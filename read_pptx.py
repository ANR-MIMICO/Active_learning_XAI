import sys
import subprocess

def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    import pptx
except ImportError:
    install('python-pptx')
    import pptx

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

if __name__ == "__main__":
    file_path = "v0.pptx"
    try:
        content = extract_text_from_pptx(file_path)
        with open("pptx_content.txt", "w", encoding="utf-8") as f:
            f.write(content)
        print("Successfully extracted pptx content to pptx_content.txt")
    except Exception as e:
        print(f"Error: {e}")
